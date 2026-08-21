"""The deck must build, and its shapes must land on the slide.

`docs/deck/build_deck.py` generates the capstone presentation. Keeping it in the
repository is only worth something if it still runs, so this builds the real
deck into a temporary directory and checks the geometry.

The bounds assertion is not theoretical: the governance table's last column
overran the right edge by 0.12in on first build, and with no renderer available
locally there was nothing else that would have caught it before the deck was in
front of an audience.

Skipped when python-pptx is absent, so a contributor who only touches the
pipeline never has to install it.
"""

import importlib.util
import math
import re
import subprocess
import sys
from pathlib import Path

import pytest

pytest.importorskip("pptx", reason="python-pptx not installed; deck tests skipped")

from pptx import Presentation  # noqa: E402 — must follow importorskip

ROOT = Path(__file__).resolve().parents[1]
BUILDER = ROOT / "docs" / "deck" / "build_deck.py"
EMU_PER_INCH = 914400
# Pinned so a slide silently disappearing from the builder is caught. Update
# deliberately when the deck genuinely changes length.
EXPECTED_SLIDES = 19


@pytest.fixture(scope="module")
def deck(tmp_path_factory):
    out = tmp_path_factory.mktemp("deck") / "deck.pptx"
    result = subprocess.run(
        [sys.executable, str(BUILDER), "--out", str(out)],
        capture_output=True, text=True, cwd=ROOT,
    )
    assert result.returncode == 0, result.stderr[-2000:]
    assert out.exists(), "builder reported success but wrote no file"
    return Presentation(str(out))


def test_deck_builds_with_every_slide(deck):
    assert len(deck.slides._sldIdLst) == EXPECTED_SLIDES


def test_every_shape_lands_on_the_slide(deck):
    """A shape past the edge is invisible in Slides and silent in the file."""
    width, height = deck.slide_width, deck.slide_height
    overflow = []
    for i, slide in enumerate(deck.slides, 1):
        for shape in slide.shapes:
            if shape.left is None:
                continue
            right, bottom = shape.left + shape.width, shape.top + shape.height
            if shape.left < 0 or shape.top < 0 or right > width or bottom > height:
                overflow.append(
                    f"slide {i}: {shape.name} "
                    f"(right={right / EMU_PER_INCH:.2f}in, bottom={bottom / EMU_PER_INCH:.2f}in)"
                )
    assert not overflow, "shapes outside the slide:\n  " + "\n  ".join(overflow)


def test_no_slide_is_empty(deck):
    for i, slide in enumerate(deck.slides, 1):
        text = "".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
        assert text.strip(), f"slide {i} has no text"


def test_text_is_unlikely_to_overflow_its_box(deck):
    """Rough estimate — catches a paragraph that grew well past its container."""
    problems = []
    for i, slide in enumerate(deck.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame or not shape.text_frame.text.strip():
                continue
            w_in = shape.width / EMU_PER_INCH
            h_in = shape.height / EMU_PER_INCH
            needed = 0.0
            for para in shape.text_frame.paragraphs:
                if not para.runs:
                    continue
                size = max((r.font.size.pt if r.font.size else 14) for r in para.runs)
                text = "".join(r.text for r in para.runs)
                per_line = max(8, int(w_in * 96 / (size * 0.52)))
                lines = max(1, math.ceil(len(text) / per_line)) + text.count("\v")
                needed += lines * size * 1.22 / 72 + (para.space_after.pt if para.space_after else 0) / 72
            if needed > h_in * 1.6 and needed - h_in > 0.5:
                problems.append(f"slide {i}: needs ~{needed:.2f}in in a {h_in:.2f}in box "
                                f"| {shape.text_frame.text[:40]!r}")
    assert not problems, "text likely overflows:\n  " + "\n  ".join(problems)


def test_screenshot_assets_are_in_the_repository():
    """A deck that depends on ~/Desktop is a deck only one machine can build."""
    assets = ROOT / "docs" / "deck" / "assets"
    referenced = set(re.findall(r'screenshot\(s,\s*"([^"]+)"', BUILDER.read_text()))
    assert referenced, "no screenshots referenced — remove this test if that is intended"
    for name in referenced:
        assert (assets / name).exists(), f"deck references {name}, which is not in docs/deck/assets/"


def test_downstream_work_is_labelled_as_a_prototype():
    """The copilot is a separate project and an early one.

    Showing it is fair; letting a slide imply it is a finished capstone
    deliverable is not. This fails if the qualifier is edited away.
    """
    source = BUILDER.read_text()
    if "copilot-overview.png" not in source:
        return
    assert "prototype" in source.lower(), (
        "the copilot screenshots are in the deck without being labelled a prototype"
    )


def test_deck_claims_match_the_repository():
    """Numbers on the slides must match what the repo actually contains.

    A deck is the easiest artefact to leave stale — it is regenerated rarely and
    read by the people least able to check it.
    """
    source = BUILDER.read_text()
    marts = {p.stem for p in (ROOT / "src" / "pipeline").glob("*.py")}
    assert marts, "no pipeline files found"

    # The deck must not state a dashboard count that contradicts the repository.
    # Asserting the *correct* phrasing is brittle — "1 AI/BI dashboards" is not
    # even grammatical — so this asserts the absence of every wrong count instead,
    # in digits and in words.
    actual = len(list((ROOT / "dashboards").glob("*.lvdash.json")))
    words = {1: "one", 2: "two", 3: "three", 4: "four", 5: "five"}
    for n, word in words.items():
        if n == actual:
            continue
        for claim in (f"{n} AI/BI dashboard", f"{word} dashboards", f"{n} dashboards"):
            assert claim not in source, (
                f"deck claims {claim!r} but the repository has {actual} dashboard(s)"
            )

    genie = ROOT / "genie" / "f1_gold_space.json"
    if genie.exists():
        assert "not yet configured" not in source, \
            "deck still describes Genie as unconfigured, but genie/ exists"
