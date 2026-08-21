"""Generate the capstone presentation. The deck is code, not a binary.

    pip install -r requirements-dev.txt
    python3 docs/deck/build_deck.py                       # writes build/F1-Race-Intelligence-Capstone.pptx
    python3 docs/deck/build_deck.py --out ~/Desktop/deck.pptx

Drag the result into Google Drive and open it — Slides converts it natively.

WHY A SCRIPT
------------
Because the deck changes whenever the platform does, and hand-editing eighteen
slides is where stale claims survive. When the Genie agent went from "chosen,
not configured" to built, updating this file was five string replacements and a
rebuild — after which `grep` could prove the old wording was gone from every
slide. That is not a check you can run against a hand-made .pptx.

It also makes the layout testable: `tests/test_deck.py` builds the deck and
asserts every shape lands inside the slide. That caught the governance table
running 0.12in off the right edge, which nothing but a renderer would otherwise
have shown.

STRUCTURE
---------
Palette and geometry constants, then layout helpers — `page` for the kicker and
title furniture, `statgrid` for the big-number cards, `points` for bold-head
bullets, `card` for panels — then `node`/`arrow_r`/`arrow_d`, which draw the
architecture diagram on slide 4 as real positioned shapes. Everything after that
is one labelled block per slide.

Changing the palette is four lines at the top. Changing the story is the blocks.
"""

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# Formula 1's own palette: the near-black #15151E its broadcast graphics sit on,
# and #E10600, the red used on the logo and every trackside board. The accent
# colours are team liveries rather than arbitrary picks — Mercedes petronas for
# a verified number, McLaren papaya for a gap, a racing blue for context — so a
# colour on a slide means something to anyone who watches the sport.
BG     = RGBColor(0x15, 0x15, 0x1E)   # F1 dark
CARD   = RGBColor(0x1F, 0x1F, 0x2B)
CARD2  = RGBColor(0x27, 0x27, 0x3A)
LINE   = RGBColor(0x3A, 0x3A, 0x50)
TEXT   = RGBColor(0xFF, 0xFF, 0xFF)
MUTED  = RGBColor(0x9A, 0x9A, 0xB0)
RED    = RGBColor(0xE1, 0x06, 0x00)   # official F1 red
TEAL   = RGBColor(0x00, 0xD2, 0xBE)   # Mercedes petronas
AMBER  = RGBColor(0xFF, 0x87, 0x00)   # McLaren papaya
BLUE   = RGBColor(0x00, 0x90, 0xFF)   # racing blue

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.85)
BW = W - 2 * M

# Titillium Web is Formula 1's own typeface and is available in Google Slides,
# where this deck is presented from. Anywhere it is missing the fallback is a
# plain sans — a different look, never a broken one.
FONT = "Titillium Web"
FONT_FALLBACK = "Arial"

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def tf(slide, l, t, w, h):
    box = slide.shapes.add_textbox(l, t, w, h)
    f = box.text_frame
    f.word_wrap = True
    f.margin_left = f.margin_right = f.margin_top = f.margin_bottom = 0
    return f


def para(f, text, size, color=TEXT, bold=False, after=8, before=0,
         first=False, align=PP_ALIGN.LEFT, space=None, face=None):
    p = f.paragraphs[0] if first else f.add_paragraph()
    p.text = text
    p.alignment = align
    p.space_after = Pt(after)
    p.space_before = Pt(before)
    if space:
        p.line_spacing = space
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        # Display sizes carry the F1 face; body copy stays on the universally
        # available fallback, so a missing font can never cost legibility.
        r.font.name = face or (FONT if size >= 18 else FONT_FALLBACK)
    return p


def card(slide, l, t, w, h, fill=CARD, border=LINE, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    if radius:
        try:
            shape.adjustments[0] = 0.08
        except Exception:
            pass
    shape.text_frame.word_wrap = True
    return shape


def dot(slide, l, t, color, size=Inches(0.12)):
    d = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, size, size)
    d.fill.solid()
    d.fill.fore_color.rgb = color
    d.line.fill.background()
    d.shadow.inherit = False
    return d


def chequer(slide, left, top, cols=6, rows=2, size=Inches(0.17), color=None):
    """A chequered-flag block. Two rows is enough to read as one; more is noise."""
    color = color or TEXT
    for r in range(rows):
        for c in range(cols):
            if (r + c) % 2:
                continue
            sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        left + c * size, top + r * size, size, size)
            sq.fill.solid()
            sq.fill.fore_color.rgb = color
            sq.line.fill.background()
            sq.shadow.inherit = False


def speed_lines(slide, right, top, color=None):
    """Three tapering bars — the motion mark F1 uses on everything."""
    color = color or RED
    for i, (w, h, gap) in enumerate([(Inches(1.5), Inches(0.09), 0),
                                     (Inches(1.1), Inches(0.09), Inches(0.18)),
                                     (Inches(0.7), Inches(0.09), Inches(0.36))]):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right - w, top + gap, w, h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        bar.shadow.inherit = False


def base(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def page(kicker, title, n=None):
    """A section slide. The footer number is the slide's real position.

    It used to be an argument, which made inserting a slide a renumbering
    exercise across every call below it — and a footer that disagrees with the
    slide it sits on is the kind of error nobody catches while rehearsing.
    Counting a private counter was no better: the first inserted slide that did
    not call page() put every footer after it out by one. Asking the deck how
    many slides it holds cannot drift.
    """
    s = prs.slides.add_slide(BLANK)
    n = len(prs.slides._sldIdLst)
    base(s)
    k = tf(s, M, Inches(0.52), BW, Inches(0.3))
    para(k, kicker.upper(), 11, RED, bold=True, first=True, after=0)
    t = tf(s, M, Inches(0.86), BW, Inches(0.72))
    para(t, title, 31, TEXT, bold=True, first=True, after=0)
    # A hairline across the page, with a short red lead-in — the start-line feel
    # without a graphic that competes with the content.
    rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, Inches(1.62), BW, Emu(9525))
    rule.fill.solid(); rule.fill.fore_color.rgb = LINE
    rule.line.fill.background(); rule.shadow.inherit = False
    lead_in = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, Inches(1.605), Inches(1.1), Inches(0.035))
    lead_in.fill.solid(); lead_in.fill.fore_color.rgb = RED
    lead_in.line.fill.background(); lead_in.shadow.inherit = False
    if n:
        f = tf(s, W - M - Inches(0.8), H - Inches(0.6), Inches(0.8), Inches(0.3))
        para(f, f"{n:02d}", 11, MUTED, first=True, align=PP_ALIGN.RIGHT, after=0)
    return s, Inches(1.95)


def lead(s, top, headline, sub=None, size=21):
    f = tf(s, M, top, BW, Inches(1.1))
    para(f, headline, size, TEXT, bold=True, first=True, after=6)
    if sub:
        para(f, sub, 14.5, MUTED, after=0, space=1.25)
    return f


def points(s, top, items, size=15, gap=15, width=None, head_color=TEXT):
    f = tf(s, M, top, width or BW, H - top - Inches(0.9))
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            para(f, item[0], size, head_color, bold=True, first=(i == 0),
                 after=3, before=0 if i == 0 else 11)
            para(f, item[1], size - 2, MUTED, after=gap, space=1.22)
        else:
            para(f, item, size, TEXT, first=(i == 0), after=gap, space=1.22)
    return f


def statgrid(s, top, rows, cols=4, h=Inches(1.35)):
    gap = Inches(0.24)
    w = (BW - gap * (cols - 1)) / cols
    for i, (value, label, color) in enumerate(rows):
        r, c = divmod(i, cols)
        l = M + c * (w + gap)
        t = top + r * (h + gap)
        card(s, l, t, w, h)
        f = tf(s, l + Inches(0.24), t + Inches(0.2), w - Inches(0.48), h - Inches(0.4))
        para(f, value, 30, color, bold=True, first=True, after=4)
        para(f, label, 11, MUTED, after=0, space=1.15)


# ══════════════════════ 1 · title ══════════════════════
s = prs.slides.add_slide(BLANK)
base(s)
glow = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.22), H)
glow.fill.solid(); glow.fill.fore_color.rgb = RED
glow.line.fill.background(); glow.shadow.inherit = False
chequer(s, Inches(11.0), Inches(0.75), cols=8, rows=2, size=Inches(0.2))
speed_lines(s, W - Inches(0.9), Inches(6.35))
f = tf(s, Inches(1.15), Inches(2.15), Inches(10.8), Inches(2.8))
para(f, "DATA ENGINEERING CAPSTONE", 12, RED, bold=True, first=True, after=14)
para(f, "Formula 1", 44, TEXT, bold=True, after=2)
para(f, "Race Intelligence & Strategy Platform", 44, MUTED, bold=True, after=20)
para(f, "A governed batch lakehouse on Databricks Free Edition.\n"
        "Public APIs → Unity Catalog Volume → Lakeflow Declarative Pipeline → AI/BI.",
     15, MUTED, after=0, space=1.3)
chips = ["3 seasons", "59 rounds", "6 Gold marts", "224 tests", "0 reconciliation mismatches"]
x = Inches(1.15)
for i, c in enumerate(chips):
    w = Inches(0.55 + 0.105 * len(c))
    ch = card(s, x, Inches(5.75), w, Inches(0.42), fill=CARD2)
    cf = ch.text_frame
    cf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(cf, c, 11, TEAL if i == 4 else MUTED, first=True, after=0, align=PP_ALIGN.CENTER)
    x += w + Inches(0.16)

# ══════════════════════ 2 · audience ══════════════════════
s, top = page("Why it exists", "Built for people who argue about races")
lead(s, top, "The audience drives the architecture more than the data volume does.")
points(s, Inches(2.9), [
    ("Reporters — checking a claim before publishing",
     "Was that win earned on pace, or inherited in the pit lane? Did the stewards change the result after the flag?"),
    ("Fans and analysts — arguing about who was actually quick",
     "The points table cannot answer it. 150 driver-races here finished at least three places behind their own pace."),
])
band = card(s, M, Inches(5.5), BW, Inches(1.15), fill=CARD2)
f = tf(s, M + Inches(0.3), Inches(5.68), BW - Inches(0.6), Inches(0.9))
para(f, "So it is sized for trust, not scale.", 17, TEAL, bold=True, first=True, after=5)
para(f, "~1,200 race results is nothing to a lakehouse. The work is making every number traceable to a raw "
        "API payload — and a wrong number findable rather than merely absent.", 13, MUTED, after=0, space=1.2)

# ══════════════════════ 3 · what is built ══════════════════════
s, top = page("Scope", "What is built")
statgrid(s, top, [
    ("11", "Bronze streaming tables\none per endpoint", TEXT),
    ("8 + 3", "Silver facts + dimensions\ntwo dimensions SCD Type 2", TEXT),
    ("6 + 1", "Gold marts, plus a governed\nmetric view over them", TEXT),
    ("6 + 1", "dashboard pages, one per decision\nplus a Genie agent over Gold", TEXT),
    ("65,862", "lap timings — the finest grain", BLUE),
    ("9", "quarantine views, one per fact", AMBER),
    ("185", "automated tests", TEAL),
    ("0", "credentials required — both APIs keyless", TEAL),
])
f = tf(s, M, Inches(5.95), BW, Inches(0.8))
para(f, "Sources: Jolpica-F1 (results · qualifying · standings · sprint · pit stops · laps) and the "
        "Open-Meteo ERA5 archive for measured race-day weather — two sources, "
        "different shapes, one medallion path.", 13, MUTED, first=True, after=0)

# ══════════════════════ 4 · architecture diagram ══════════════════════
s, top = page("Architecture", "One path, checked at every boundary")

def node(l, t, w, h, title, sub, accent=None, fill=CARD):
    c = card(s, l, t, w, h, fill=fill)
    if accent:
        strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, Inches(0.055), h)
        strip.fill.solid(); strip.fill.fore_color.rgb = accent
        strip.line.fill.background(); strip.shadow.inherit = False
    f = tf(s, l + Inches(0.22), t + Inches(0.16), w - Inches(0.36), h - Inches(0.3))
    para(f, title, 13.5, TEXT, bold=True, first=True, after=3)
    para(f, sub, 10.5, MUTED, after=0, space=1.12)
    return c

def arrow_r(l, t, w=Inches(0.34), h=Inches(0.22)):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = LINE
    a.line.fill.background(); a.shadow.inherit = False

def arrow_d(l, t, w=Inches(0.22), h=Inches(0.3)):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, l, t, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = LINE
    a.line.fill.background(); a.shadow.inherit = False

NW, NH = Inches(3.35), Inches(0.95)
GAP = Inches(0.5)
y1 = Inches(2.0)
node(M, y1, NW, NH, "Sources",
     "Jolpica-F1 REST  ·  Open-Meteo ERA5\npublic, keyless", RED)
arrow_r(M + NW + Inches(0.08), y1 + Inches(0.37))
node(M + NW + GAP, y1, NW, NH, "Ingestion  ·  Databricks Job",
     "serverless Python, throttled, retried,\nidempotent partition writes", RED)
arrow_r(M + 2 * NW + GAP + Inches(0.08), y1 + Inches(0.37))
node(M + 2 * (NW + GAP), y1, NW, NH, "Landing  ·  UC Volume",
     "f1.raw.landing — raw JSON per call,\nprovenance envelope", RED)

arrow_d(M + 2 * (NW + GAP) + NW / 2, y1 + NH + Inches(0.06))

y2 = Inches(3.42)
PW = Inches(2.62)
PG = Inches(0.36)
pipe_x = M
node(pipe_x, y2, PW, Inches(1.15), "Bronze",
     "Auto Loader streaming tables\nread as text · nothing dropped", BLUE)
arrow_r(pipe_x + PW + Inches(0.02), y2 + Inches(0.46))
node(pipe_x + PW + PG, y2, PW, Inches(1.15), "Silver",
     "facts: MV + natural-key dedupe\ndims: Auto CDC SCD Type 2", BLUE)
arrow_r(pipe_x + 2 * PW + PG + Inches(0.02), y2 + Inches(0.46))
node(pipe_x + 2 * (PW + PG), y2, PW, Inches(1.15), "Gold",
     "6 marts, dims joined as-of race\ndate + governed metric view", BLUE)
arrow_r(pipe_x + 3 * PW + 2 * PG + Inches(0.02), y2 + Inches(0.46))
node(pipe_x + 3 * (PW + PG), y2, PW, Inches(1.15), "Serving",
     "Dashboard: 6 decision pages\nMetric view · Genie agent", TEAL)

y3 = Inches(4.78)
sw = (BW - Inches(0.36)) / 2
side = card(s, M, y3, sw, Inches(0.62), fill=CARD2)
f = tf(s, M + Inches(0.2), y3 + Inches(0.16), sw - Inches(0.4), Inches(0.4))
para(f, "Rejected rows → 9 quarantine tables, with the rule they violated", 11.5, AMBER, first=True, after=0)
side2 = card(s, M + sw + Inches(0.36), y3, sw, Inches(0.62), fill=CARD2)
f = tf(s, M + sw + Inches(0.56), y3 + Inches(0.16), sw - Inches(0.4), Inches(0.4))
para(f, "Pipeline event log → f1.gold.pipeline_event_log, queryable in SQL", 11.5, AMBER, first=True, after=0)

y4 = Inches(5.62)
gov = card(s, M, y4, BW, Inches(0.62), fill=CARD)
strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, y4, Inches(0.055), Inches(0.62))
strip.fill.solid(); strip.fill.fore_color.rgb = TEAL
strip.line.fill.background(); strip.shadow.inherit = False
f = tf(s, M + Inches(0.24), y4 + Inches(0.16), BW - Inches(0.5), Inches(0.4))
para(f, "Unity Catalog  —  governs all of it: lineage from the pipeline graph, layered grants, audit, "
        "row-level provenance on every Bronze row", 12, TEXT, first=True, after=0)
f = tf(s, M, Inches(6.45), BW, Inches(0.5))
para(f, "Batch, single path — not Lambda, not Kappa. F1 produces data 24 times a year, in bursts, hours after each race.",
     12, MUTED, first=True, after=0)

# ══════════════════════ 5 · the detailed diagram ══════════════════════
# Full bleed and no slide chrome: the picture carries its own title, and every
# inch given to a heading is an inch taken off text that is already small.
# Generated by docs/deck/build_architecture.py — run that first.
s = prs.slides.add_slide(BLANK)
base(s)
_arch = Path(__file__).resolve().parents[1] / "assets" / "architecture.png"
assert _arch.exists(), (
    f"missing {_arch} — run: python3 docs/deck/build_architecture.py"
)
_img_w = W
_img_h = int(W / 2.069)          # the diagram's aspect
s.shapes.add_picture(str(_arch), 0, int((H - _img_h) / 2), width=int(_img_w), height=_img_h)

# ══════════════════════ 6 · technology choices ══════════════════════
s, top = page("Stack", "Technology choices, and why")
tech = [
    ("Databricks Free Edition", "Free, serverless, a real Lakehouse. The goal of the project is Databricks depth.", TEAL),
    ("Python & SQL", "API ingestion in Python; transformations and marts in SQL.", TEAL),
    ("Unity Catalog Volume", "Managed, governed file landing zone — no external location needed. The Free-Edition-correct choice.", TEAL),
    ("Databricks Job (scheduled)", "Serverless ingestion. Replaces AWS Lambda + EventBridge: one platform, no cloud account.", TEAL),
    ("Auto Loader", "Incremental file processing and checkpointing without an always-on stream.", TEAL),
    ("Lakeflow Declarative Pipeline", "Declarative medallion, built-in expectations, lineage. Triggered, so it runs batch.", TEAL),
    ("Delta Lake", "Reliable lakehouse table format with time travel.", TEAL),
    ("Unity Catalog", "Lineage, access control, schema management, audit.", TEAL),
    ("Unity Catalog metric views", "The semantic layer. One definition of a point, shared by the dashboard and Genie.", TEAL),
    ("AI/BI Dashboards", "Reporting over Gold — one dashboard, a page per analyst decision.", TEAL),
    ("Genie", "Natural-language analytics, scoped to the six Gold marts and the metric view above them. 7 certified queries.", TEAL),
]
row_h = Inches(0.44)
y = top - Inches(0.05)
for i, (name, why, mark) in enumerate(tech):
    if i % 2 == 0:
        bg = card(s, M, y, BW, row_h, fill=CARD, border=BG)
    dot(s, M + Inches(0.2), y + Inches(0.16), mark, Inches(0.1))
    f = tf(s, M + Inches(0.45), y + Inches(0.1), Inches(3.5), row_h)
    para(f, name, 12.5, TEXT, bold=True, first=True, after=0)
    f = tf(s, M + Inches(4.1), y + Inches(0.1), BW - Inches(4.3), row_h)
    para(f, why, 12, MUTED, first=True, after=0)
    y += row_h
f = tf(s, M, Inches(6.65), BW, Inches(0.4))
para(f, "Every row is built and running. Ingestion is Python; everything downstream is declarative.",
     11, MUTED, first=True, after=0)

# ══════════════════════ 6 · dataset types ══════════════════════
s, top = page("Design", "The dataset type is the architecture")
lead(s, top, "Not a style preference.",
     "Pick the wrong one and it shows up as stale aggregates, or a full recompute nobody asked for.")
points(s, Inches(3.0), [
    ("Bronze — streaming table + Auto Loader",
     "Files arrive incrementally and are never revised in place. Read as text: optional Ergast fields appear in only some files, and schema inference would fail the pipeline the first time an unusual one landed."),
    ("Silver facts — materialised view, deduplicated on the natural key",
     "A result is provisional when published; stewards reorder it after the flag. Keeping the greatest _ingest_ts is a full-partition window function, which streaming append mode cannot express."),
    ("Silver dimensions — streaming table + Auto CDC, SCD Type 2",
     "A driver's constructor changes and that history is the point. Auto CDC needs a streaming source, so these read Bronze directly — a materialised view cannot be streamed from."),
    ("Gold — materialised view, batch read",
     "MVs recompute when an upstream row is amended. Streaming tables would not, and amendment is exactly the case that matters."),
], size=14, gap=5)

# ══════════════════════ 7 · decisions ══════════════════════
s, top = page("Decisions", "Three that were expensive to discover")
points(s, top, [
    ("SCD-2 is built from results, not from /drivers",
     "Every field /drivers returns is static, so Auto CDC over it yields a dimension with no history — the pattern implemented but never exercised. The attribute that changes is the constructor, and it only appears in results.  →  42 versions across 28 drivers, 14 historical."),
    ("Sprint points are part of the championship",
     "Summing race points alone leaves 13 of 24 drivers short of their official 2024 total. With sprint points, reconciliation against the independent standings endpoint is exact."),
    ("Idempotency is a two-part contract, and both halves are required",
     "Ingestion skips closed rounds and re-pulls the open one; Silver deduplicates by natural key on the newest _ingest_ts. Drop either half and the live round double-counts."),
], size=15, gap=10)

# ══════════════════════ 8 · data quality ══════════════════════
s, top = page("Data quality", "Rejected rows are routed, never dropped on the floor")
points(s, top, [
    "Rules are declared once as dicts and used twice: expect_all_or_drop on the fact, and the inverse on a quarantine table that records which rules the row violated.",
    "A static pre-flight checks every column named in an expectation exists in the staged view — before a cluster starts, because graph analysis costs quota to learn about a typo.",
], size=14, gap=10)
lead(s, Inches(3.35), "The census is not zero, and should not be", size=17)
statgrid(s, Inches(4.05), [
    ("69", "lap rows outside 40–300 s\nred-flag and safety-car laps", AMBER),
    ("8", "standings rows with no\nchampionship position", AMBER),
    ("2", "pit stops published with\nan empty duration", AMBER),
    ("0", "rows lost without a\nreason attached", TEAL),
], h=Inches(1.25))
f = tf(s, M, Inches(5.65), BW, Inches(0.6))
para(f, "Zero everywhere would mean the expectations had stopped being evaluated.", 13, MUTED, first=True, after=0)

# ══════════════════════ 9 · correctness ══════════════════════
s, top = page("Correctness", "Proven, not asserted")
lead(s, top, "A green pipeline update is not the bar. Reconciliation is.",
     "The strongest check reconciles two independent endpoints against each other: points earned from results + sprint, "
     "versus the published driverStandings endpoint. Agreement is evidence the flattening, deduplication and joins are "
     "correct — not a comparison against a number typed in from memory.")
statgrid(s, Inches(3.65), [
    ("0", "reconciliation mismatches\nevery driver, every season", TEAL),
    ("10 / 10", "validation checks passing\nlast task of both jobs", TEAL),
    ("14", "SCD-2 history rows\nthe dimension is truly versioned", TEAL),
    ("0", "duplicate natural keys\nanywhere in Silver", TEAL),
])
f = tf(s, M, Inches(5.4), BW, Inches(1.0))
para(f, "Lap timings reconcile against results with two exemptions, both real semantics rather than tolerance for error: "
        "disqualification zeroes the results lap count while the laps driven were still timed, and a lapped runner is "
        "classified on the lap the leader finished.", 13, MUTED, first=True, after=0, space=1.2)

# ══════════════════════ 10 · governance ══════════════════════
s, top = page("Governance", "Access narrows as the data gets rawer")
rows = [("Principal", "Catalog", "Gold", "Silver / Bronze", "Landing Volume"),
        ("account users", "USE_CATALOG", "SELECT", "—", "—"),
        ("engineer tier", "USE_CATALOG", "SELECT", "SELECT", "READ_VOLUME"),
        ("owner", "everything by ownership", "", "", "")]
rh = Inches(0.5)
for r, row in enumerate(rows):
    y = top + r * rh
    if r == 0:
        card(s, M, y, BW, rh, fill=CARD2, border=BG)
    elif r % 2 == 1:
        card(s, M, y, BW, rh, fill=CARD, border=BG)
    xs = [M + Inches(0.25), M + Inches(2.95), M + Inches(5.15), M + Inches(6.95), M + Inches(9.25)]
    widths = [Inches(2.6), Inches(2.1), Inches(1.7), Inches(2.2), Inches(2.2)]
    for c, val in enumerate(row):
        f = tf(s, xs[c], y + Inches(0.14), widths[c], Inches(0.35))
        col = TEXT if r == 0 else (TEAL if val == "SELECT" and c == 2 else MUTED)
        para(f, val, 12, col, bold=(r == 0), first=True, after=0)
points(s, Inches(4.15), [
    "An analyst can query f1.gold.driver_performance and cannot read the Bronze payload behind it. Applied by an idempotent script through the permissions API — no compute, works with the daily quota exhausted.",
    "Nobody is granted WRITE_VOLUME. The landing zone has exactly one writer, the ingestion job; a second breaks the idempotency contract the deduplication depends on.",
    ("Honest limit",
     "On a single-user workspace the owner's ownership outranks every grant, so the layering is real in configuration but cannot be experienced. Proving it needs a second identity that owns nothing."),
], size=13, gap=8)

# ══════════════════════ 11 · semantic layer ══════════════════════
s, top = page("Semantic layer", "One definition of a point")
lead(s, top, "A metric used to be defined once per consumer.",
     "total_points as a mart column, pace against team in a dashboard dataset, and \"always use "
     "total_points\" written into the Genie agent's instructions. Those agreed because one person "
     "wrote all three — which is not a guarantee, it is a coincidence with good intentions.")
points(s, Inches(3.35), [
    ("f1.gold.driver_metrics — a Unity Catalog metric view",
     "13 measures across 9 dimensions, queried through MEASURE(). Team resolves to the as-of-race "
     "constructor, so the historical attribution survives into the semantic layer instead of being "
     "re-derived by whoever asks next."),
    ("The dashboard and Genie both read it",
     "Migrated row for row: zero rows disagreeing on starts, points, qualifying, pace, pit stop or "
     "DNF rate. The definitions moved; the numbers did not."),
    ("Points Per Start is why the mechanism matters",
     "A ratio that re-aggregates correctly at any grain — driver, team, season, circuit. A standard "
     "view cannot do that, because its aggregation is fixed the moment it is created."),
], size=14, gap=6)

# ══════════════════════ 12 · the dashboard ══════════════════════
s, top = page("The product", "Six pages, one per decision")
lead(s, top, "Organised by the question being asked, not by the table underneath.",
     "Pages named after marts are an implementation detail leaking into the product.")
rows = [
    ("Standings", "Where did the season finish?", "points · wins · gap · title margin"),
    ("Driver Form", "Who is over-performing their car?", "pace vs own team · places made up · consistency"),
    ("Constructor Benchmarking", "Where are points won and lost?", "quali · race pace · pit stop · DNF rate"),
    ("Circuit Priors", "What usually works here?", "typical stops · overtaking index · pole-to-win · weather"),
    ("Championship Swing", "Which rounds decided it?", "gap to leader · points gained per round"),
    ("Trust", "Can these numbers be believed?", "expectations passed and failed, per dataset"),
]
rh = Inches(0.62)
for i, (name, question, metrics) in enumerate(rows):
    y = Inches(3.2) + i * rh
    if i % 2 == 0:
        card(s, M, y, BW, rh, fill=CARD, border=BG)
    f = tf(s, M + Inches(0.25), y + Inches(0.1), Inches(3.1), rh)
    para(f, name, 13, TEXT, bold=True, first=True, after=0)
    f = tf(s, M + Inches(3.5), y + Inches(0.1), Inches(4.0), rh)
    para(f, question, 12.5, TEAL, first=True, after=0)
    f = tf(s, M + Inches(7.7), y + Inches(0.12), BW - Inches(7.9), rh)
    para(f, metrics, 11.5, MUTED, first=True, after=0)
f = tf(s, M, Inches(7.0), BW, Inches(0.35))
para(f, "Plus a Genie agent over the same Gold layer, for the questions no tile anticipated.",
     12, MUTED, first=True, after=0)

# ══════════════════════ 13 · automation ══════════════════════
s, top = page("Operations", "Automated, bounded, and watched")
points(s, top, [
    ("f1_end_to_end  —  on demand",
     "unit tests → ingest → pipeline → validation. Tests first, so a typo costs seconds instead of a cluster start."),
    ("f1_ingest_incremental  —  weekly, paused in dev",
     "ingest → pipeline → validation. Validation sits on the scheduled path, because correctness checked only when a human remembers is a habit, not a guarantee."),
    ("Both bounded and alerting",
     "One-hour timeout, failure email, retries only on the idempotent task. On Free Edition a hung task spends tomorrow's compute allowance as well as today's — silently."),
    ("Portable by construction",
     "No workspace host, profile or season is hardcoded anywhere. A fork runs ./scripts/bootstrap.sh --profile <name> and gets the same platform."),
], size=14.5, gap=6)

# ══════════════════════ 14 · testing ══════════════════════
s, top = page("Assurance", "Testing in three layers")
points(s, top, [
    ("224 local tests  ·  no Spark  ·  under a second",
     "Ingestion logic — pagination against a total that counts inner records, the rate budget, the closed-round predicate, season derivation — plus repository contracts: no legacy DLT API, correct dashboard widget versions, no catalog fallback, every job bounded, every scheduled job validating."),
    ("Static pre-flight across all five Silver files",
     "Every column named in an expectation must exist in the staged view. Otherwise found only by graph analysis: a cluster start and quota, to learn about a typo."),
    ("Spark tests on Databricks, first task of the job",
     "The real lap-time and pit-duration parsers and the three-level laps schema, loaded from the actual pipeline files with a stubbed decorator module — so no transformation logic is duplicated into the test."),
], size=14, gap=8)
f = tf(s, M, Inches(6.1), BW, Inches(0.6))
para(f, "Every contract assertion corresponds to a mistake actually made here that cost a failed update or a blank dashboard tile.",
     12.5, MUTED, first=True, after=0)

# ══════════════════════ 15 · findings ══════════════════════
s, top = page("Findings", "What the data actually says")
points(s, top, [
    ("Pace and result are different questions",
     "Ranking each driver's median clean lap within a race against where they finished: 150 driver-races ended at least three places behind their own pace, 232 ahead of it. No points table shows that."),
    ("Rain does not cause chaos — and the dashboard argues the point",
     "Monza 2024 measured 19.1 mm and ran dry. A daily total cannot tell rain that fell overnight from rain that fell during the race, so the mart names the source of the flag rather than asserting it rained. The wet-versus-dry bars on Circuit Priors would separate if rain changed the racing. They barely do."),
    ("Running the field's strategy is worth about 1.1 positions",
     "Stay with the field and you gain 0.69 places on average; deviate in either direction and you lose 0.40. Stints are derived from stops, because nothing publishes them — and 'two stops' means nothing until you know the field made three."),
], size=15, gap=10)

# ══════════════════════ 16 · future scope ══════════════════════
s, top = page("Next", "From dashboards to a copilot")
lead(s, top, "The marts are the product. The next step is who else can reach them.",
     "Everything below reads the same Gold layer — no second copy of the data, no new pipeline.")
cards = [
    ("Fan & reporter app", "Databricks Apps front end over the Genie agent: ask in plain English, get the number and the query behind it.", BLUE),
    ("Semantic search", "Vector Search over race reports and regulations, so 'why was he penalised' returns the text, not a row.", BLUE),
    ("Standings & weather explorer", "Every season's championship progression and measured race-day conditions, on demand.", BLUE),
    ("Amendment history", "Answering what the stewards changed needs the Silver facts as streaming tables — a real trade, not a checkbox.", AMBER),
]
gap = Inches(0.24)
cw = (BW - gap * 3) / 4
for i, (title, sub, col) in enumerate(cards):
    l = M + i * (cw + gap)
    card(s, l, Inches(3.7), cw, Inches(1.85))
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, Inches(3.7), cw, Inches(0.05))
    strip.fill.solid(); strip.fill.fore_color.rgb = col
    strip.line.fill.background(); strip.shadow.inherit = False
    f = tf(s, l + Inches(0.22), Inches(3.95), cw - Inches(0.44), Inches(1.5))
    para(f, title, 13.5, TEXT, bold=True, first=True, after=5)
    para(f, sub, 11, MUTED, after=0, space=1.18)
f = tf(s, M, Inches(5.85), BW, Inches(0.8))
para(f, "The Genie agent already answers questions over Gold today, governed by the same grants. The app is the front "
        "door for people who will never open a Databricks workspace. SCD-2 already gives version history on the "
        "dimensions — 42 driver versions, queryable now.", 13, TEAL, first=True, after=0, space=1.2)

# ══════════════════════ 17 · close ══════════════════════
s = prs.slides.add_slide(BLANK)
base(s)
glow = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.22), H)
glow.fill.solid(); glow.fill.fore_color.rgb = RED
glow.line.fill.background(); glow.shadow.inherit = False
chequer(s, Inches(11.0), Inches(0.75), cols=8, rows=2, size=Inches(0.2))
f = tf(s, Inches(1.15), Inches(2.3), Inches(10.8), Inches(2.4))
para(f, "IN ONE LINE", 12, RED, bold=True, first=True, after=16)
para(f, "A pipeline whose numbers", 34, TEXT, bold=True, after=2)
para(f, "you can check —", 34, TEXT, bold=True, after=2)
para(f, "and whose gaps you can read.", 34, MUTED, bold=True, after=18)
para(f, "Ingestion is idempotent, Silver deduplicates on natural keys, dimensions carry history, Gold joins them "
        "as of the race date, and ten validation checks say so on every run — including the weekly one nobody watches.",
     14, MUTED, after=0, space=1.3)
statgrid(s, Inches(5.6), [
    ("3", "seasons · 59 rounds", MUTED),
    ("6", "Gold marts · 6 decision pages", MUTED),
    ("224", "automated tests", MUTED),
    ("0", "reconciliation mismatches", TEAL),
], h=Inches(1.15))

_parser = argparse.ArgumentParser(description="Build the capstone deck.")
_parser.add_argument(
    "--out",
    default=str(Path(__file__).resolve().parents[2] / "build" / "F1-Race-Intelligence-Capstone.pptx"),
    help="output path for the .pptx (default: build/ at the repository root)",
)
_args = _parser.parse_args()
_out = Path(_args.out).expanduser()
_out.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(_out))
print(f"wrote {_out}  ({len(prs.slides._sldIdLst)} slides)")
